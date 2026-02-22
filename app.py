import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import io

st.set_page_config(page_title="PDF a PPTX Profesional", page_icon="🎨")

st.title("🎨 Convertidor con Diseño y Texto Editable")
st.write("Mantiene el diseño original y añade texto editable encima.")

uploaded_file = st.file_uploader("Sube el PDF de NotebookLM", type="pdf")

if uploaded_file is not None:
    if st.button("Generar PowerPoint"):
        with st.spinner("Procesando diseño y texto..."):
            prs = Presentation()
            # Configurar tamaño de diapositiva panorámica (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            pdf_data = uploaded_file.read()
            doc = fitz.open(stream=pdf_data, filetype="pdf")
            
            for page in doc:
                # 1. Convertir página a imagen para el fondo
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_data = pix.tobytes("png")
                
                # 2. Crear diapositiva
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 3. Insertar la imagen de fondo
                img_stream = io.BytesIO(img_data)
                slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                
                # 4. OCR para detectar texto y ponerlo ENCIMA
                img_pil = Image.open(io.BytesIO(img_data))
                d = pytesseract.image_to_data(img_pil, output_type=pytesseract.Output.DICT, lang='spa')
                
                # Agrupamos por bloques para no tener mil cuadritos
                for i in range(len(d['text'])):
                    if int(d['conf'][i]) > 50:
                        text = d['text'][i].strip()
                        if len(text) > 1:
                            # Ajustar coordenadas al tamaño de la diapositiva
                            w_img, h_img = img_pil.size
                            left = Inches((d['left'][i] / w_img) * 13.33)
                            top = Inches((d['top'][i] / h_img) * 7.5)
                            width = Inches((d['width'][i] / w_img) * 13.33)
                            height = Inches((d['height'][i] / h_img) * 7.5)
                            
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.text_frame
                            p = tf.paragraphs[0]
                            p.text = text
                            p.font.size = Pt(11)
                            # Nota: El texto se verá un poco encimado al original, 
                            # pero ahora puedes borrar el original o editar este.

            pptx_io = io.BytesIO()
            prs.save(pptx_io)
            pptx_io.seek(0)
            
            st.success("¡Generado con éxito!")
            st.download_button("📥 Descargar PPTX", pptx_io, "final.pptx")
